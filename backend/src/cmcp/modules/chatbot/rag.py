from __future__ import annotations

import hashlib
import os
import re
import uuid
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Any, Optional
from urllib.parse import unquote, urlparse

from sqlalchemy import select

from cmcp.config.database import db
from cmcp.config.media_config import settings as media_settings
from cmcp.config.settings import settings
from cmcp.core.exceptions import BusinessValidationError
from cmcp.modules.academic.models import Course, CourseChapter, CourseOffering, Semester
from cmcp.modules.chatbot.models import ChatMessage, ChatbotMaterialIndex
from cmcp.modules.materials.models import Material
from cmcp.modules.media.encrypted_files import download_and_decrypt_file
from cmcp.modules.media.storage import get_backend


RE_CHAPTER = re.compile(
    r"\b(?:chapter|ch\.?|section)\s*(\d+|one|two|three|four|five|six|seven|eight|nine|ten)\b",
    re.IGNORECASE,
)
RE_QNA = re.compile(
    r"\b(generate[\s_]questions?[\s_]and[\s_]answers?|questions?[\s_]and[\s_]answers?|q\s*(?:&|and)\s*a)\b",
    re.IGNORECASE,
)
RE_QUIZ = re.compile(
    r"\b(quiz+(?:e|z)?|generate[\s_]quiz+(?:e|z)?|make[\s_]quiz+(?:e|z)?|test[\s_]me|mcq|practice[\s_](?:test|questions?))\b",
    re.IGNORECASE,
)
RE_SUMMARY = re.compile(
    r"\b(summar(?:y|ize|ise|ized?)|overview|recap|brief|chapter[\s_]summary)\b",
    re.IGNORECASE,
)
RE_JUNK = re.compile(
    r"(ppt/slides/slide\d+\.xml|slideLayout\d+\.xml|slideMaster\d+\.xml|notesSlide\d+\.xml|<\?xml|xmlns:|<a:t>|<p:sp>|\.rels|fntdata|\.png)",
    re.IGNORECASE,
)


def compute_hash(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def normalize_label(value: str) -> str:
    return str(value or "").strip().lower()


def semester_label(semester: Optional[Semester]) -> str:
    if not semester:
        return "Unassigned Semester"
    if semester.name:
        return semester.name.strip()
    return f"Semester {semester.number}"


def material_filename(material: Material) -> str:
    file_url = material.file_url or ""
    parsed = urlparse(file_url)
    name = os.path.basename(parsed.path or file_url).replace(".enc", "")
    if name and "." in name:
        return name
    title = material.title or f"material-{material.id}"
    ext = ""
    if "." in name:
        ext = "." + name.rsplit(".", 1)[-1]
    return f"{title}{ext}"


def _extension_from_name(filename: str) -> str:
    clean = filename[:-4] if filename.endswith(".enc") else filename
    return clean.rsplit(".", 1)[-1].lower() if "." in clean else ""


def _extract_text(file_bytes: bytes, filename: str) -> str:
    ext = _extension_from_name(filename)

    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(file_bytes))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext in {"ppt", "pptx"}:
        from pptx import Presentation

        prs = Presentation(BytesIO(file_bytes))
        slides = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            lines = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(("- " if para.level > 0 else "") + line)
            if len(lines) > 1:
                slides.append("\n".join(lines))
        return "\n\n".join(slides)

    if ext == "docx":
        from docx import Document

        doc = Document(BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)

    return file_bytes.decode("utf-8", errors="ignore")


def _is_junk_chunk(text: str) -> bool:
    if RE_JUNK.search(text):
        return True
    words = text.split()
    if not words:
        return True
    path_like = sum(1 for word in words if "/" in word or word.endswith(".xml") or word.endswith(".rels"))
    return path_like / len(words) > 0.3


def _chunk_text(text: str, chunk_size: int = 400, overlap: int = 80) -> list[str]:
    words = text.split()
    step = max(1, chunk_size - overlap)
    chunks = []
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk and not _is_junk_chunk(chunk):
            chunks.append(chunk)
    return chunks


def _chapter_label_from_filename(filename: str) -> Optional[str]:
    stem = filename.rsplit(".", 1)[0].replace("_", " ").strip()
    if re.match(r"^(chapter|section)\s*\d+", stem, re.IGNORECASE):
        return stem.title()
    return None


def _chapter_label_from_text(text_chunk: str) -> Optional[str]:
    for line in text_chunk.split("\n")[:3]:
        if re.match(r"^\s*(chapter|section)\s*\d+", line, re.IGNORECASE):
            return line.strip()
    return None


def _resolve_chapter(text_chunk: str, filename: str, chapter: Optional[CourseChapter]) -> str:
    if chapter:
        return f"Chapter {chapter.number}: {chapter.title}"
    return _chapter_label_from_text(text_chunk) or _chapter_label_from_filename(filename) or "General Content"


def _file_key_from_url(file_url: str) -> Optional[str]:
    if not file_url:
        return None
    marker = "/api/media/file/"
    if marker in file_url:
        return unquote(file_url.split(marker, 1)[1])
    parsed = urlparse(file_url)
    path = parsed.path or file_url
    local_base = media_settings.LOCAL_PUBLIC_BASE.rstrip("/") + "/"
    if path.startswith(local_base):
        return unquote(path[len(local_base):])
    return None


def read_material_file(material: Material) -> tuple[bytes, str]:
    file_url = material.file_url or ""
    if not file_url:
        raise BusinessValidationError("Material has no file attached.")

    key = _file_key_from_url(file_url)
    filename = material_filename(material)
    if key:
        if key.endswith(".enc"):
            data, _mime = download_and_decrypt_file(key)
            return data, filename or os.path.basename(key).replace(".enc", "")
        return get_backend().download(key), filename or os.path.basename(key)

    parsed = urlparse(file_url)
    if parsed.scheme in {"", "file"}:
        path = Path(parsed.path if parsed.scheme == "file" else file_url)
        if path.exists() and path.is_file():
            return path.read_bytes(), filename or path.name

    raise BusinessValidationError("Cannot read this material file for indexing.")


class RagProvider:
    def __init__(self):
        self._collection = None
        self._llm = None

    @property
    def collection(self):
        if self._collection is None:
            try:
                import chromadb
            except Exception as exc:
                raise BusinessValidationError("chromadb is not installed.") from exc
            chroma_dir = Path(settings.CHATBOT_CHROMA_DIR)
            if not chroma_dir.is_absolute():
                chroma_dir = Path(__file__).resolve().parents[4] / settings.CHATBOT_CHROMA_DIR
            chroma_dir.mkdir(parents=True, exist_ok=True)
            client = chromadb.PersistentClient(path=str(chroma_dir))
            self._collection = client.get_or_create_collection(name=settings.CHATBOT_COLLECTION_NAME)
        return self._collection

    @property
    def llm(self):
        if self._llm is None:
            if not settings.DEEPSEEK_API_KEY:
                raise BusinessValidationError("DEEPSEEK_API_KEY is not configured.")
            try:
                from openai import OpenAI
            except Exception as exc:
                raise BusinessValidationError("openai is not installed.") from exc
            self._llm = OpenAI(api_key=settings.DEEPSEEK_API_KEY, base_url=settings.CHATBOT_LLM_BASE_URL)
        return self._llm

    def call_llm(self, system_prompt: str, user_prompt: str) -> str:
        resp = self.llm.chat.completions.create(
            model=settings.CHATBOT_LLM_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        )
        return resp.choices[0].message.content


@lru_cache(maxsize=1)
def provider() -> RagProvider:
    return RagProvider()


def delete_material_chunks(company_id: int, material_id: int) -> int:
    coll = provider().collection
    res = coll.get(where={"$and": [{"company_id": int(company_id)}, {"material_id": int(material_id)}]}, include=[])
    ids = res.get("ids") or []
    if ids:
        coll.delete(ids=ids)
    return len(ids)


def index_material(material: Material, *, force: bool = False) -> dict[str, Any]:
    file_bytes, filename = read_material_file(material)
    file_hash = compute_hash(file_bytes)

    existing = ChatbotMaterialIndex.query.filter_by(
        company_id=int(material.company_id),
        material_id=int(material.id),
    ).first()

    if existing and existing.file_hash == file_hash and not force:
        return {"material_id": int(material.id), "status": "skipped", "chunk_count": int(existing.chunk_count)}

    text = _extract_text(file_bytes, filename)
    if not text.strip():
        raise BusinessValidationError("No readable text could be extracted from this material.")

    chunks = _chunk_text(text)
    if not chunks:
        raise BusinessValidationError("No useful text chunks could be extracted from this material.")

    delete_material_chunks(int(material.company_id), int(material.id))

    offering = material.course_offering
    subject = offering.course.title if offering and offering.course else material.title
    sem_label = semester_label(offering.semester if offering else None)
    norm_sem = normalize_label(sem_label)
    norm_subject = subject.strip()
    chapter = material.chapter

    ids = [f"mat_{material.id}_{uuid.uuid4().hex[:12]}" for _ in chunks]
    metadatas = [
        {
            "company_id": int(material.company_id),
            "material_id": int(material.id),
            "semester": norm_sem,
            "subject": norm_subject,
            "source": filename,
            "chapter": _resolve_chapter(chunk, filename, chapter),
            "chunk_index": idx,
        }
        for idx, chunk in enumerate(chunks)
    ]
    provider().collection.add(documents=chunks, ids=ids, metadatas=metadatas)

    if not existing:
        existing = ChatbotMaterialIndex(
            company_id=int(material.company_id),
            material_id=int(material.id),
            file_hash=file_hash,
            chunk_count=len(chunks),
            semester_label=norm_sem,
            subject_name=norm_subject,
            source_name=filename,
        )
        db.session.add(existing)
    else:
        existing.file_hash = file_hash
        existing.chunk_count = len(chunks)
        existing.semester_label = norm_sem
        existing.subject_name = norm_subject
        existing.source_name = filename

    db.session.flush()
    return {"material_id": int(material.id), "status": "indexed", "chunk_count": len(chunks)}


def _chat_history(session_id: str, last_n: int = 6) -> str:
    rows = (
        ChatMessage.query
        .filter_by(session_id=session_id)
        .order_by(ChatMessage.created_at.asc())
        .all()[-last_n:]
    )
    return "\n".join(f"{r.role_name.upper()}: {r.message_text}" for r in rows)


def _vector_query(query_text: str, n: int, where: Optional[dict] = None) -> tuple[list, list, list]:
    kwargs = {
        "query_texts": [query_text],
        "n_results": n,
        "include": ["documents", "metadatas", "distances"],
    }
    if where:
        kwargs["where"] = where
    res = provider().collection.query(**kwargs)
    return (
        res.get("documents", [[]])[0],
        res.get("metadatas", [[]])[0],
        res.get("distances", [[]])[0],
    )


def _filter_relevant(docs: list, metas: list, distances: list) -> list[tuple[str, dict]]:
    return [
        (doc, meta)
        for doc, meta, distance in zip(docs, metas, distances)
        if distance < settings.CHATBOT_RELEVANCE_THRESHOLD
    ]


def _detect_chapter(question: str) -> Optional[str]:
    match = RE_CHAPTER.search(question)
    return match.group(0).strip() if match else None


def _subject_chapters(company_id: int, norm_semester: str, subject_name: str) -> list[str]:
    res = provider().collection.get(
        where={"$and": [
            {"company_id": int(company_id)},
            {"semester": norm_semester},
            {"subject": subject_name},
        ]},
        include=["metadatas"],
    )
    chapters = {
        meta.get("chapter")
        for meta in (res.get("metadatas") or [])
        if meta.get("chapter") and meta.get("chapter") != "General Content"
    }
    return sorted(chapters) or ["General Content"]


def _chapter_chunks(company_id: int, norm_semester: str, subject_name: str, chapter_label: str, limit: int = 30) -> list[str]:
    res = provider().collection.get(
        where={"$and": [
            {"company_id": int(company_id)},
            {"semester": norm_semester},
            {"subject": subject_name},
            {"chapter": chapter_label},
        ]},
        include=["documents"],
        limit=limit,
    )
    docs = res.get("documents") or []
    if docs:
        return docs

    all_res = provider().collection.get(
        where={"$and": [
            {"company_id": int(company_id)},
            {"semester": norm_semester},
            {"subject": subject_name},
        ]},
        include=["documents", "metadatas"],
        limit=500,
    )
    needle = chapter_label.lower()
    return [
        doc for doc, meta in zip(all_res.get("documents") or [], all_res.get("metadatas") or [])
        if needle in (meta.get("chapter") or "").lower() or needle in (meta.get("source") or "").lower()
    ][:limit]


def _special_content_answer(company_id: int, norm_semester: str, subject_name: str, question: str, mode: str, chapter_label: Optional[str]) -> dict[str, Any]:
    available = _subject_chapters(company_id, norm_semester, subject_name)
    if not chapter_label:
        options = "\n".join(f"- {chapter}" for chapter in available)
        label = {"quiz": "Quiz Mode", "qna": "Q&A Generation Mode", "summary": "Summary Mode"}[mode]
        return {
            "answer": f"**{label}**\n\nWhich chapter should I use?\n\n{options}\n\nReply with a chapter, for example: `{mode} chapter 1`.",
            "sources": [],
            "mode": mode,
            "chapters": available,
        }

    matched = next((chapter for chapter in available if chapter_label.lower() in chapter.lower()), None)
    display = matched or chapter_label
    chunks = _chapter_chunks(company_id, norm_semester, subject_name, display)
    if not chunks:
        return {"answer": f"No content found for **{display}** in **{subject_name}**.", "sources": []}

    context = "\n\n".join(chunks[:settings.CHATBOT_MAX_CONTEXT_CHUNKS])
    system = f"You are a strict academic tutor for {subject_name}. Use only the provided material."
    if mode == "quiz":
        prompt = f"""Generate a 10-question multiple-choice quiz from this content.

Use this exact format:
**Q1.** Question text

A) Option
B) Option
C) Option
D) Option

After all questions add:
---ANSWER KEY---
Q1: A

Content:
{context}"""
    elif mode == "qna":
        prompt = f"""Generate 10 open-ended questions with full answers from this content.

Use this exact format:
**Q1.** Question text

**Answer:** Detailed answer.

Content:
{context}"""
    else:
        prompt = f"""Write a structured chapter summary from this content.

Include:
## Overview
## Key Concepts
## Key Takeaways

Content:
{context}"""

    return {
        "answer": provider().call_llm(system, prompt),
        "sources": [{"chapter": display, "subject": subject_name}],
        "mode": mode,
    }


def answer_question(company_id: int, session_id: str, semester: str, subject_name: str, question: str) -> dict[str, Any]:
    norm_semester = normalize_label(semester)
    norm_subject = subject_name.strip()
    chapter = _detect_chapter(question)
    is_qna = bool(RE_QNA.search(question))
    is_quiz = bool(RE_QUIZ.search(question)) and not is_qna
    is_summary = bool(RE_SUMMARY.search(question))

    if is_qna:
        return _special_content_answer(company_id, norm_semester, norm_subject, question, "qna", chapter)
    if is_quiz:
        return _special_content_answer(company_id, norm_semester, norm_subject, question, "quiz", chapter)
    if is_summary:
        return _special_content_answer(company_id, norm_semester, norm_subject, question, "summary", chapter)

    docs, metas, dists = _vector_query(
        question,
        n=settings.CHATBOT_TOP_K,
        where={"$and": [
            {"company_id": int(company_id)},
            {"semester": norm_semester},
            {"subject": norm_subject},
        ]},
    )
    relevant = _filter_relevant(docs, metas, dists)
    if not relevant:
        return {
            "answer": "I do not have enough indexed material to answer that. Please ask an admin to index the course materials.",
            "sources": [],
        }

    context = "\n\n".join(
        f"[Source: {meta.get('source', '?')} | {meta.get('chapter', 'General')}]\n{doc}"
        for doc, meta in relevant[:settings.CHATBOT_MAX_CONTEXT_CHUNKS]
    )
    history = _chat_history(session_id)
    system = (
        f"You are a strict academic tutor for {subject_name} ({semester}). "
        "Answer only from the provided context. Never use outside knowledge."
    )
    user = f"""Context:
{context}

Chat History:
{history}

Question: {question}

Rules:
- Use Markdown.
- Keep paragraphs short.
- Cite sources inline like [Source: filename.pdf].
- If the answer is not in the context, say the material does not cover it."""

    return {
        "answer": provider().call_llm(system, user),
        "sources": [meta for _, meta in relevant],
    }
