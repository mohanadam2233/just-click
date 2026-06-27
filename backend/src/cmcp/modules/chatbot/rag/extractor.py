from __future__ import annotations

import hashlib
import os
import re
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote, urlparse

from cmcp.config.media_config import settings as media_settings
from cmcp.core.exceptions import BusinessValidationError
from cmcp.modules.materials.models import Material
from cmcp.modules.media.encrypted_files import download_and_decrypt_file
from cmcp.modules.media.storage import get_backend


RE_JUNK = re.compile(
    r"(ppt/slides/slide\d+\.xml|slideLayout\d+\.xml|slideMaster\d+\.xml|notesSlide\d+\.xml|<\?xml|xmlns:|<a:t>|<p:sp>|\.rels|fntdata|\.png)",
    re.IGNORECASE,
)
RE_PRINTABLE_ASCII = re.compile(rb"[\x20-\x7E\r\n\t]{4,}")
RE_PRINTABLE_UNICODE = re.compile(rb"(?:[\x20-\x7E\r\n\t]\x00){4,}")
RE_PPT_NOISE = re.compile(
    r"^(Picture|Word\.Picture|Rectangle \d+|Title Placeholder|Text Placeholder|"
    r"Date Placeholder|Footer Placeholder|Slide Number Placeholder|___PPT\d+|"
    r"Century Schoolbook|Arial|Wingdings|Times New Roman|Calibri|Helvetica|SimHei|"
    r"Courier New|Monotype Sorts|sohne|hne Mono|View)$",
    re.IGNORECASE,
)


def compute_hash(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def material_filename(material: Material) -> str:
    file_url = material.file_url or ""
    parsed = urlparse(file_url)
    name = os.path.basename(parsed.path or file_url).replace(".enc", "")
    if name and "." in name:
        return name
    title = material.title or f"material-{material.id}"
    ext = ""
    if "." in name:
        ext = "." + name.rsplit(".", 1)[-1]
    return f"{title}{ext}"


def _extension_from_name(filename: str) -> str:
    clean = filename[:-4] if filename.endswith(".enc") else filename
    return clean.rsplit(".", 1)[-1].lower() if "." in clean else ""


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    lines.append(("- " if para.level > 0 else "") + line)
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_ppt_legacy(file_bytes: bytes) -> str:
    """Extract text from legacy PowerPoint 97-2003 (.ppt) OLE files."""
    try:
        import olefile
    except ImportError as exc:
        raise BusinessValidationError(
            "Legacy .ppt extraction requires olefile. Run: pip install olefile"
        ) from exc

    ole = olefile.OleFileIO(BytesIO(file_bytes))
    try:
        if not ole.exists("PowerPoint Document"):
            raise BusinessValidationError("Legacy .ppt file is missing PowerPoint Document stream.")

        data = ole.openstream("PowerPoint Document").read()
    finally:
        ole.close()

    chunks: list[str] = []
    for match in RE_PRINTABLE_ASCII.finditer(data):
        text = match.group(0).decode("ascii", errors="ignore").strip()
        if len(text) >= 4 and not RE_PPT_NOISE.match(text):
            chunks.append(text)

    for match in RE_PRINTABLE_UNICODE.finditer(data):
        text = match.group(0).decode("utf-16-le", errors="ignore").strip()
        if len(text) >= 4 and not RE_PPT_NOISE.match(text):
            chunks.append(text)

    deduped: list[str] = []
    seen: set[str] = set()
    for chunk in chunks:
        normalized = " ".join(chunk.split())
        if normalized in seen:
            continue
        seen.add(normalized)
        deduped.append(normalized)

    return "\n".join(deduped)


def _extract_text(file_bytes: bytes, filename: str) -> str:
    ext = _extension_from_name(filename)

    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(file_bytes))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == "pptx":
        return _extract_pptx(file_bytes)

    if ext == "ppt":
        return _extract_ppt_legacy(file_bytes)

    if ext == "docx":
        from docx import Document

        doc = Document(BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)

    return file_bytes.decode("utf-8", errors="ignore")


def _file_key_from_url(file_url: str) -> str | None:
    if not file_url:
        return None
    marker = "/api/media/file/"
    if marker in file_url:
        return unquote(file_url.split(marker, 1)[1])
    parsed = urlparse(file_url)
    path = parsed.path or file_url
    local_base = media_settings.LOCAL_PUBLIC_BASE.rstrip("/") + "/"
    if path.startswith(local_base):
        return unquote(path[len(local_base):])
    return None


def read_material_file(material: Material) -> tuple[bytes, str]:
    file_url = material.file_url or ""
    if not file_url:
        raise BusinessValidationError("Material has no file attached.")

    key = _file_key_from_url(file_url)
    filename = material_filename(material)
    if key:
        if key.endswith(".enc"):
            data, _mime = download_and_decrypt_file(key)
            return data, filename or os.path.basename(key).replace(".enc", "")
        return get_backend().download(key), filename or os.path.basename(key)

    parsed = urlparse(file_url)
    if parsed.scheme in {"", "file"}:
        path = Path(parsed.path if parsed.scheme == "file" else file_url)
        if path.exists() and path.is_file():
            return path.read_bytes(), filename or path.name

    raise BusinessValidationError("Cannot read this material file for indexing.")
